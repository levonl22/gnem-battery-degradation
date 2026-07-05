"""Build the 12-slide final deck → results/gnem_battery_slides.pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "results/figures"
OUT = ROOT / "results/gnem_battery_slides.pptx"


def add_title_slide(prs, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs, title: str, bullets: list[str], fig=None, fig2=None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
    top = Inches(1.35)
    if fig and fig.exists():
        left = Inches(0.5) if fig2 else Inches(1.5)
        width = Inches(5.8) if fig2 else Inches(10)
        slide.shapes.add_picture(str(fig), left, top, width=width)
    if fig2 and fig2.exists():
        slide.shapes.add_picture(str(fig2), Inches(6.8), top, width=Inches(5.8))


def add_image_slide(prs, title: str, bullets: list[str], fig: Path, width=Inches(10.5)) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(16)
    if fig.exists():
        slide.shapes.add_picture(str(fig), Inches(1.2), Inches(1.55), width=width)


def main() -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "Early Prediction of Lithium-Ion Battery End-of-Life",
        "GNEM Research Fellow · Levon Lau · July 2026\n"
        "github.com/levonl22/gnem-battery-degradation",
    )
    add_bullets(
        prs,
        "Problem & motivation",
        [
            "Full life tests run hundreds to thousands of cycles — predict EOL from early data",
            "EOL: first cycle with discharge capacity < 80% of initial (cycle_index ≥ 1)",
            "Dataset: MIT-Stanford-Toyota fast-charging (Severson et al., 2019)",
            "Pipeline: tabular ML vs GRU sequence model → monotonic SOH → window ablation",
        ],
    )
    add_bullets(
        prs,
        "Dataset & pipeline",
        [
            "140 raw JSON files → 134 unique cells (dedupe; one partial cell dropped)",
            "about 111k cycle rows · 44 cell-level features",
            "Cell-level split: 94 train / 20 val / 20 test (random_state = 42)",
            "Notebooks 01–10 · processed CSVs · metrics JSON · figures in results/",
        ],
    )
    add_bullets(
        prs,
        "Capacity fade & EOL distribution",
        ["134 cells · median EOL about 792 cycles · range 159–2,237"],
        FIG / "capacity_fade_example.png",
        FIG / "eol_distribution.png",
    )
    add_image_slide(
        prs,
        "Early-cycle features",
        [
            "44 features: capacity, SOH, resistance, efficiency, temperature + ΔV(Q)",
            "Strongest correlates: ΔV(Q) spread (|r| about 0.8) · efficiency (r about 0.78)",
        ],
        FIG / "feature_correlation.png",
        width=Inches(9.5),
    )
    add_image_slide(
        prs,
        "ML baselines — test error (20 cells)",
        [
            "Linear 133 · ElasticNet 132 · Random forest 101 · "
            "XGBoost 85 cycles MAE (about 11% MAPE)",
        ],
        FIG / "model_comparison_baselines.png",
    )
    add_image_slide(
        prs,
        "ML — predicted vs true EOL",
        [
            "XGBoost best on holdout; ΔV(Q) and efficiency drive tree importance",
            "20 test cells → indicative metrics",
        ],
        FIG / "pred_vs_true_eol_baselines.png",
    )
    add_image_slide(
        prs,
        "GRU sequence model",
        [
            "Input: 100 cycles × 4 channels (SOH, resistance, efficiency, temperature)",
            "Test MAE about 110 cycles — behind XGBoost (about 85); no ΔV(Q) in sequence",
        ],
        FIG / "pred_vs_true_eol_sequence.png",
    )
    add_image_slide(
        prs,
        "Monotonic SOH constraint",
        [
            "Dual-head GRU + monotonic penalty · EOL MAE about 112 cycles "
            "(unchanged vs unconstrained)",
        ],
        FIG / "soh_curves_constrained.png",
        width=Inches(8.5),
    )
    add_image_slide(
        prs,
        "Early-cycle ablation (N = 20, 50, 100)",
        [
            "XGBoost: 167 → 107 → 85  |  GRU: 144 → 143 → 110 cycles test MAE",
            "Best model depends on window length and available features",
        ],
        FIG / "ablation_early_cycles.png",
    )
    add_bullets(
        prs,
        "Limitations",
        [
            "134 cells, 20-cell test holdout → high metric variance",
            "Single LFP/graphite fast-charge protocol; transfer untested",
            "XGBoost uses ΔV(Q) from raw JSON; GRU uses four summary channels only",
            "Mechanism-agnostic features; fixed 80% EOL definition",
        ],
    )
    add_bullets(
        prs,
        "Conclusions & future work",
        [
            "Best model: XGBoost at N = 100 — about 85 cycles test MAE (about 11% MAPE)",
            "About 50–100 early cycles + voltage-curve features support screening",
            "Future: NASA/CALCE validation · ΔV(Q) in sequences · larger cohorts",
            "Report + repo: github.com/levonl22/gnem-battery-degradation",
        ],
    )
    prs.save(str(OUT))
    print(f"Wrote {len(prs.slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
