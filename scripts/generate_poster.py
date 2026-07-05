"""Fill GNEM showcase poster template → results/gnem_battery_poster.pptx."""

import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "results/figures"
OUT = ROOT / "results/gnem_battery_poster.pptx"
TEMPLATE = Path.home() / "Downloads/GNEM_Project_Poster_Template.pptx"

# Showcase layout: heatmap = star visual (Chart 44, largest chart well).
CONTENT: dict[str, str] = {
    "TextBox 5": (
        "Early Prediction of Lithium-Ion Battery End-of-Life\n"
        "from Early-Cycle Data"
    ),
    "TextBox 6": (
        "Levon Lau · GNEM Summer Fellow · University of Georgia · July 2026"
    ),
    "TextBox 15": "Why this matters",
    "TextBox 16": (
        "• Full battery life tests take hundreds to thousands of cycles — slow and expensive\n"
        "• Goal: predict end-of-life (EOL) from early measurements before obvious fade\n"
        "• EOL = first cycle below 80% of initial discharge capacity\n"
        "• Dataset: MIT-Stanford-Toyota fast-charging cells (Severson et al., 2019)"
    ),
    "TextBox 21": "Approach",
    "TextBox 17": (
        "• 140 raw files → 134 unique cells · cell-level 94 / 20 / 20 split\n"
        "• 44 hand-crafted early-cycle features + Severson-style ΔV(Q) curves\n"
        "• Classical ML (XGBoost, random forest, …) vs PyTorch GRU on cycle trajectories\n"
        "• Ablation: how many of the first N cycles are needed? (N = 20, 50, 100)"
    ),
    "TextBox 19": "Headline results",
    "TextBox 20": (
        "• Best model: XGBoost at N = 100 — about 85 cycles test MAE (about 11% MAPE)\n"
        "• GRU sequence model: about 110 cycles MAE (competitive, not best)\n"
        "• More early cycles help XGBoost most: 167 → 107 → 85 cycles MAE\n"
        "• ΔV(Q) and energy efficiency matter more than raw capacity early on"
    ),
    "TextBox 42": "Feature signals vs lifetime (exploratory)",
    "TextBox 43": (
        "Heatmap (center): top 12 features most correlated with EOL across 134 cells.\n"
        "Red = higher feature values linked to longer life · blue = shorter life.\n"
        "Voltage-curve shape (ΔV(Q)) and efficiency stand out before capacity visibly fades."
    ),
    "TextBox 22": "Takeaways",
    "TextBox 26": (
        "• Early-cycle data can estimate lifetime — useful for fast-charge cell screening\n"
        "• About 50–100 cycles + voltage-curve features give the strongest tabular signal\n"
        "• Reproducible open pipeline on GitHub (notebooks 01–10)"
    ),
    "TextBox 27": (
        "Limitations: 20-cell test holdout; single chemistry/protocol; 134 cells total.\n"
        "Future: NASA/CALCE validation · richer sequence inputs · larger cohorts.\n"
        "github.com/levonl22/gnem-battery-degradation"
    ),
}

# Chart 44 = largest chart slot → heatmap hero
CHART_IMAGES = {
    "Chart 44": FIG / "feature_correlation_poster.png",
    "Chart 46": FIG / "ablation_early_cycles.png",
    "Chart 45": FIG / "model_comparison_baselines.png",
}
HERO_IMAGE = FIG / "capacity_fade_example.png"

FONT_SIZES = {
    "TextBox 5": 48,
    "TextBox 6": 20,
    "TextBox 15": 26,
    "TextBox 19": 26,
    "TextBox 21": 26,
    "TextBox 22": 26,
    "TextBox 42": 24,
}


def set_text(shape, text: str, font_size: int) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)


def replace_shape_with_image(slide, shape_name: str, image_path: Path) -> None:
    for shape in slide.shapes:
        if shape.name != shape_name:
            continue
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        shape._element.getparent().remove(shape._element)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return


def main() -> None:
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")

    subprocess.run(
        [sys.executable, str(ROOT / "scripts/generate_poster_heatmap.py")],
        check=True,
    )

    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    for shape in slide.shapes:
        if shape.name not in CONTENT:
            continue
        set_text(shape, CONTENT[shape.name], FONT_SIZES.get(shape.name, 15))

    for chart_name, img in CHART_IMAGES.items():
        if img.exists():
            replace_shape_with_image(slide, chart_name, img)
        else:
            print(f"Warning: missing {img}")

    if HERO_IMAGE.exists():
        replace_shape_with_image(slide, "Picture 4", HERO_IMAGE)

    prs.save(str(OUT))
    print(f"Wrote → {OUT}")


if __name__ == "__main__":
    main()
